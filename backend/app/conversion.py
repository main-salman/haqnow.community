"""
Document conversion service to standardize all documents to PDF format
"""

import io
import logging
import os
import subprocess
import tempfile
from pathlib import Path
from typing import BinaryIO, Optional, Tuple

import fitz  # PyMuPDF
import pandas as pd
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentConverter:
    """Converts various document formats to PDF"""

    @staticmethod
    def convert_to_pdf(file_data: bytes, filename: str) -> Tuple[bytes, str]:
        """
        Convert any supported document type to PDF

        Args:
            file_data: Raw file bytes
            filename: Original filename for format detection

        Returns:
            Tuple of (pdf_bytes, converted_filename)
        """
        file_ext = Path(filename).suffix.lower()
        base_name = Path(filename).stem

        logger.info(f"Converting {filename} ({file_ext}) to PDF")

        # If already PDF, return as-is
        if file_ext == ".pdf":
            return file_data, filename

        # Route to appropriate conversion method
        if file_ext in [".docx", ".doc"]:
            return DocumentConverter._convert_word_to_pdf(file_data, base_name)
        elif file_ext in [".pptx", ".ppt"]:
            return DocumentConverter._convert_powerpoint_to_pdf(file_data, base_name)
        elif file_ext in [".xlsx", ".xls", ".csv"]:
            return DocumentConverter._convert_spreadsheet_to_pdf(
                file_data, base_name, file_ext
            )
        elif file_ext in [".txt"]:
            return DocumentConverter._convert_text_to_pdf(file_data, base_name)
        elif file_ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"]:
            return DocumentConverter._convert_image_to_pdf(file_data, base_name)
        else:
            # Try LibreOffice conversion as fallback
            return DocumentConverter._convert_with_libreoffice(file_data, filename)

    @staticmethod
    def _convert_word_to_pdf(file_data: bytes, base_name: str) -> Tuple[bytes, str]:
        """Convert Word document to PDF using python-docx and PyMuPDF"""
        try:
            # Try LibreOffice first for best quality
            return DocumentConverter._convert_with_libreoffice(
                file_data, f"{base_name}.docx"
            )
        except Exception as e:
            logger.warning(
                f"LibreOffice conversion failed, trying docx extraction: {e}"
            )

            # Fallback: Extract text and create PDF manually
            try:
                doc = DocxDocument(io.BytesIO(file_data))
                text_content = ""
                for paragraph in doc.paragraphs:
                    text_content += paragraph.text + "\\n"

                if not text_content.strip():
                    text_content = (
                        f"[Converted from {base_name}.docx - content extraction failed]"
                    )

                return DocumentConverter._create_pdf_from_text(text_content, base_name)
            except Exception as inner_e:
                logger.error(f"Word document conversion failed: {inner_e}")
                raise ValueError(f"Unable to convert Word document: {inner_e}")

    @staticmethod
    def _convert_powerpoint_to_pdf(
        file_data: bytes, base_name: str
    ) -> Tuple[bytes, str]:
        """Convert PowerPoint to PDF"""
        try:
            # Try LibreOffice first
            return DocumentConverter._convert_with_libreoffice(
                file_data, f"{base_name}.pptx"
            )
        except Exception as e:
            logger.warning(
                f"LibreOffice conversion failed, trying pptx extraction: {e}"
            )

            # Fallback: Extract text and create PDF
            try:
                prs = Presentation(io.BytesIO(file_data))
                text_content = ""
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_content += f"Slide {slide_num}:\\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\\n"
                    text_content += "\\n"

                if not text_content.strip():
                    text_content = (
                        f"[Converted from {base_name}.pptx - content extraction failed]"
                    )

                return DocumentConverter._create_pdf_from_text(text_content, base_name)
            except Exception as inner_e:
                logger.error(f"PowerPoint conversion failed: {inner_e}")
                raise ValueError(f"Unable to convert PowerPoint: {inner_e}")

    @staticmethod
    def _convert_spreadsheet_to_pdf(
        file_data: bytes, base_name: str, file_ext: str
    ) -> Tuple[bytes, str]:
        """Convert spreadsheet to PDF"""
        try:
            # Try LibreOffice first
            return DocumentConverter._convert_with_libreoffice(
                file_data, f"{base_name}{file_ext}"
            )
        except Exception as e:
            logger.warning(
                f"LibreOffice conversion failed, trying pandas conversion: {e}"
            )

            # Fallback: Use pandas to read and convert
            try:
                if file_ext == ".csv":
                    df = pd.read_csv(io.BytesIO(file_data))
                else:
                    df = pd.read_excel(io.BytesIO(file_data))

                # Convert DataFrame to HTML then to PDF
                html_content = df.to_html()
                return DocumentConverter._convert_html_to_pdf(html_content, base_name)
            except Exception as inner_e:
                logger.error(f"Spreadsheet conversion failed: {inner_e}")
                raise ValueError(f"Unable to convert spreadsheet: {inner_e}")

    @staticmethod
    def _convert_text_to_pdf(file_data: bytes, base_name: str) -> Tuple[bytes, str]:
        """Convert plain text to PDF"""
        try:
            text_content = file_data.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text_content = file_data.decode("latin1")
            except:
                text_content = (
                    f"[Binary content from {base_name}.txt - unable to decode as text]"
                )

        return DocumentConverter._create_pdf_from_text(text_content, base_name)

    @staticmethod
    def _convert_image_to_pdf(file_data: bytes, base_name: str) -> Tuple[bytes, str]:
        """Convert image to PDF"""
        try:
            # Use PyMuPDF to create PDF from image
            img_doc = fitz.open()
            img = fitz.open("img", file_data)
            pdf_bytes = img.convert_to_pdf()
            img.close()

            pdf_filename = f"{base_name}.pdf"
            return pdf_bytes, pdf_filename

        except Exception as e:
            logger.error(f"Image conversion failed: {e}")
            raise ValueError(f"Unable to convert image: {e}")

    @staticmethod
    def _convert_with_libreoffice(file_data: bytes, filename: str) -> Tuple[bytes, str]:
        """Convert document using LibreOffice headless mode"""
        base_name = Path(filename).stem

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file
            input_path = os.path.join(temp_dir, filename)
            with open(input_path, "wb") as f:
                f.write(file_data)

            # Convert to PDF using LibreOffice
            try:
                result = subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )

                if result.returncode != 0:
                    raise subprocess.CalledProcessError(
                        result.returncode, result.args, result.stderr
                    )

                # Read converted PDF
                pdf_path = os.path.join(temp_dir, f"{base_name}.pdf")
                if not os.path.exists(pdf_path):
                    raise FileNotFoundError(
                        f"LibreOffice did not create expected PDF: {pdf_path}"
                    )

                with open(pdf_path, "rb") as f:
                    pdf_data = f.read()

                return pdf_data, f"{base_name}.pdf"

            except subprocess.TimeoutExpired:
                raise ValueError("LibreOffice conversion timed out")
            except subprocess.CalledProcessError as e:
                raise ValueError(f"LibreOffice conversion failed: {e.stderr}")

    @staticmethod
    def _create_pdf_from_text(text_content: str, base_name: str) -> Tuple[bytes, str]:
        """Create a PDF from plain text content"""
        try:
            # Create PDF using PyMuPDF
            doc = fitz.open()
            page = doc.new_page()

            # Add text to page with basic formatting
            point = fitz.Point(50, 72)  # Start position
            font_size = 11
            line_height = font_size * 1.2

            lines = text_content.split("\\n")
            for line in lines:
                if point.y > 750:  # Near bottom of page
                    page = doc.new_page()
                    point = fitz.Point(50, 72)

                # Split long lines
                if len(line) > 80:
                    words = line.split(" ")
                    current_line = ""
                    for word in words:
                        if len(current_line + " " + word) <= 80:
                            current_line += (" " if current_line else "") + word
                        else:
                            if current_line:
                                page.insert_text(
                                    point, current_line, fontsize=font_size
                                )
                                point.y += line_height
                                if point.y > 750:
                                    page = doc.new_page()
                                    point = fitz.Point(50, 72)
                            current_line = word

                    if current_line:
                        page.insert_text(point, current_line, fontsize=font_size)
                        point.y += line_height
                else:
                    page.insert_text(point, line, fontsize=font_size)
                    point.y += line_height

            pdf_data = doc.write()
            doc.close()

            return pdf_data, f"{base_name}.pdf"

        except Exception as e:
            logger.error(f"PDF creation from text failed: {e}")
            raise ValueError(f"Unable to create PDF from text: {e}")

    @staticmethod
    def _convert_html_to_pdf(html_content: str, base_name: str) -> Tuple[bytes, str]:
        """Convert HTML content to PDF"""
        # Simple HTML to PDF conversion using PyMuPDF
        try:
            # Create a temporary HTML file and use it
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".html", delete=False
            ) as f:
                f.write(html_content)
                html_path = f.name

            try:
                # Convert HTML to PDF using PyMuPDF's story feature
                doc = fitz.open()
                story = fitz.Story(html_content)

                while not story.done:
                    page = doc.new_page()
                    rect = page.rect
                    story.place(rect)
                    story.draw(page)

                pdf_data = doc.write()
                doc.close()

                return pdf_data, f"{base_name}.pdf"
            finally:
                os.unlink(html_path)

        except Exception as e:
            logger.error(f"HTML to PDF conversion failed: {e}")
            # Fallback to text conversion
            import re
            from html import unescape

            text_content = re.sub("<[^<]+?>", "", html_content)
            text_content = unescape(text_content)
            return DocumentConverter._create_pdf_from_text(text_content, base_name)


def convert_document_to_pdf(file_data: bytes, filename: str) -> Tuple[bytes, str]:
    """
    Main function to convert any document to PDF

    Args:
        file_data: Raw file bytes
        filename: Original filename

    Returns:
        Tuple of (pdf_bytes, pdf_filename)
    """
    converter = DocumentConverter()
    return converter.convert_to_pdf(file_data, filename)
